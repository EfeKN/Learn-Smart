"""
TODO:
After migration to AWS S3, modify the save and delete methods in the FileManager 
classes to interact with the S3 bucket instead of the local file system.
"""

from io import BytesIO
import os
import pymupdf
from pptx import Presentation
from docx import Document
from fastapi import UploadFile
from abc import ABC, abstractmethod
from PIL import Image
from tools import splitext


class BaseFile(ABC):
    """
    Base class for file management.

    Attributes:
        file (UploadFile): The file to be managed.
        path (str): The path where the file is saved.

    Methods:
        save(path: str): Save the file to the specified path.
        delete(): Delete the file.
        content(): Abstract method to get the content of the file.
        get(): Abstract method to get the file.
    """

    class ResourceWrapper:
        """
        A context manager that wraps a resource and provides additional functionality.

        This class allows you to use a resource as a context manager, ensuring that the resource
        is properly closed when exiting the context. It also provides a way to access attributes
        of the resource directly.

        Usage:
        ```
        file = FileFactory()(path="file.pdf")
        with file.get() as pdf:
            # do something

        DON'T use file.resource directly as the resource may not be closed properly.
        ```

        Attributes:
            resource: The resource object to be wrapped.
        """

        def __init__(self, resource):
            self.resource = resource

        def __enter__(self):
            return self.resource

        def __exit__(self, exc_type, exc_val, exc_tb):
            if hasattr(self.resource, 'close'):
                self.resource.close()

        def __getattr__(self, attr):
            return getattr(self.resource, attr)

    def __init__(self, file: UploadFile = None, path: str = None):
        self.file = file
        self.path = path

    def save(self, path: str):
        """
        Save the file to the specified path.

        Args:
            path (str): The path where the file should be saved.

        Raises:
            OSError: If the file already exists in the specified path.
            ValueError: If no file is provided to be saved.
        """
        if self.path:
            return
        if os.path.exists(path):
            raise OSError(f"File already exists in path {path}")
        if not self.file:
            raise ValueError("No file provided to be saved.")
        
        with open(path, "wb") as file:
            file.write(self.file.file.read())

        self.path = path

    def delete(self):
        """
        Delete the file.

        If the file does not exist or no path is set, this method does nothing.
        """
        if not self.path:
            return
        if not os.path.exists(self.path):
            return
        os.remove(self.path)
        self.path = None

    @abstractmethod
    def content(self):
        """
        Abstract method to get the content of the file.

        This method should be implemented in the derived classes.
        """
        pass

    @abstractmethod
    def get(self):
        """
        Abstract method to get the file.

        This method should be implemented in the derived classes.
        """
        pass


class ImageFile(BaseFile):
    """
    Represents an image file.

    Args:
        file (UploadFile, optional): The uploaded file object. Defaults to None.
        path (str, optional): The path to the file. Defaults to None.

    Raises:
        ValueError: If the content type or extension of the file is not supported.

    Attributes:
        file (UploadFile): The uploaded file object.
        path (str): The path to the file.

    Methods:
        save: Saves the image file with optional resizing.
        content: Returns the content of the image file.
        get: Returns the image file as a ResourceWrapper object.
    """

    def __init__(self, file: UploadFile = None, path: str = None):
        ext = splitext(file.filename if file else path if path else "")[1] # extension of the file
        if ext not in ["png", "jpeg", "jpg"]:
            raise ValueError("Unsupported extension: " + ext)
        if file and file.content_type not in ["image/png", "image/jpeg", "image/jpg"]:
            raise ValueError("Unsupported content type: " + file.content_type)
        super().__init__(file, path)

    def save(self, path: str, size: tuple = (256, 256)):
        """
        Saves the image file with optional resizing.

        Args:
            path (str): The path to save the image file.
            size (tuple, optional): The desired size of the image. Defaults to (256, 256).
        """
        super().save(path)
        with Image.open(self.path) as img: # resize to the given size
            img.thumbnail(size)
            img.save(self.path)

    def content(self):
        """
        Returns the content of the image file.

        Returns:
            Image: The image object representing the content of the file.

        Raises:
            ValueError: If no file is provided.
        """
        if not self.path and not self.file:
            raise ValueError("No file provided.")
        if self.path:            
            return Image.open(self.path)
        
        return Image.open(self.file.file)
    
    def get(self):
        """
        Returns the image file as a ResourceWrapper object.

        Returns:
            ResourceWrapper: The image file wrapped in a ResourceWrapper object.
        """
        img = Image.open(self.path)
        return self.ResourceWrapper(img)
    

class PresentationFile(BaseFile):
    """
    Represents a presentation file.

    Args:
        file (UploadFile, optional): The uploaded file. Defaults to None.
        path (str, optional): The path to the file. Defaults to None.

    Raises:
        ValueError: If the content type or extension of the file is not supported.

    Attributes:
        file (UploadFile): The uploaded file.
        path (str): The path to the file.

    Methods:
        content(): Extracts the text content from the presentation file.
        get(): Retrieves the presentation resource.

    """

    def __init__(self, file: UploadFile = None, path: str = None):
        ext = splitext(file.filename if file else path if path else "")[1]
        if ext != "pptx":
            raise ValueError("Unsupported extension: " + ext)
        if file and file.content_type != "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            raise ValueError("Unsupported content type: " + file.content_type)
        super().__init__(file, path)

    def content(self):
        """
        Extracts the text content from the presentation file.

        Returns:
            str: The extracted text content.

        Raises:
            ValueError: If no file is provided.

        """
        if not self.path and not self.file:
            raise ValueError("No file provided.")
        
        text = ""

        if self.path:
            presentation = Presentation(self.path)
        else:
            presentation = Presentation(BytesIO(self.file.file.read()))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text += shape.text_frame.text
        return text
    
    def get(self):
        """
        Retrieves the presentation resource.

        Returns:
            ResourceWrapper: The presentation resource.

        """
        presentation = Presentation(self.path)
        return self.ResourceWrapper(presentation)
    

class PDFFile(BaseFile):
    """
    Represents a PDF file.

    Args:
        file (UploadFile, optional): The uploaded file object. Defaults to None.
        path (str, optional): The path to the file. Defaults to None.

    Raises:
        ValueError: If the content type or extension of the file is not supported.

    Attributes:
        file (UploadFile): The uploaded file object.
        path (str): The path to the file.

    Methods:
        content(): Returns the content of the PDF file as text.
        get(): Returns a resource wrapper for the PDF file.

    """

    def __init__(self, file: UploadFile = None, path: str = None):
        ext = splitext(file.filename if file else path if path else "")[1]
        if ext != "pdf":
            raise ValueError("Unsupported extension: " + ext)
        if file and file.content_type != "application/pdf":
            raise ValueError("Unsupported content type: " + file.content_type)
        super().__init__(file, path)

    def content(self):
        """
        Returns the content of the PDF file as text.

        Returns:
            str: The content of the PDF file.

        Raises:
            ValueError: If no file is provided.

        """
        if not self.path and not self.file:
            raise ValueError("No file provided.")
        
        # TODO: if file size is not too large and there aren't many pages, convert to images

        # TODO: Do testing with contents
        if self.path:
            with pymupdf.open(self.path) as doc:
                return chr(12).join([page.get_text() for page in doc])
        
        # TODO: find a solution for large pdf files such as books
        with pymupdf.open(stream=BytesIO(self.file.file.read()), filetype="pdf") as doc:
            return chr(12).join([page.get_text() for page in doc])

    def get(self):
        """
        Returns a resource wrapper for the PDF file.

        Returns:
            ResourceWrapper: A resource wrapper for the PDF file.

        """
        doc = pymupdf.open(self.path)
        return self.ResourceWrapper(doc)
    

class WordFile(BaseFile):
    """
    Represents a Word file.

    Args:
        file (UploadFile, optional): The uploaded file. Defaults to None.
        path (str, optional): The file path. Defaults to None.

    Raises:
        ValueError: If the content type or extension is unsupported.

    Attributes:
        file (UploadFile): The uploaded file.
        path (str): The file path.

    Methods:
        content: Returns the content of the Word file.
        get: Returns a resource wrapper for the Word file.
    """

    def __init__(self, file: UploadFile = None, path: str = None):
        ext = splitext(file.filename if file else path if path else "")[1]
        if ext != "docx":
            raise ValueError("Unsupported extension: " + ext)
        if file and file.content_type != "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            raise ValueError("Unsupported content type: " + file.content_type)
        super().__init__(file, path)

    # TODO: not tested with a docx file
    def content(self):
        """
        Returns the content of the Word file.

        Returns:
            str: The content of the Word file.
        
        Raises:
            ValueError: If no file is provided.
        """
        if not self.path and not self.file:
            raise ValueError("No file provided.")
        
        if self.path:
            doc = Document(self.path)
        else:
            doc = Document(BytesIO(self.file.file.read()))
        # TODO: Do testing
        return chr(12).join([para.text for para in doc.paragraphs])

    def get(self):
        """
        Returns a resource wrapper for the Word file.

        Returns:
            ResourceWrapper: The resource wrapper for the Word file.
        """
        doc = Document(self.path)
        return self.ResourceWrapper(doc)
    

class FileFactory:
    """
    A factory class for creating different types of files based on their extensions.
    """

    def __init__(self):
        self.file_types = {
            "pdf": PDFFile,
            "pptx": PresentationFile,
            "docx": WordFile,
            "png": ImageFile,
            "jpg": ImageFile,
            "jpeg": ImageFile
        }

    def __call__(self, file: UploadFile = None, path: str = None):
        """
        Creates an instance of a file based on the provided file or path.

        Args:
            file (UploadFile, optional): The file to be processed. Defaults to None.
            path (str, optional): The path to the file. Defaults to None.

        Returns:
            An instance of the corresponding file type based on the file extension.

        Raises:
            ValueError: If no file or path is provided.
            ValueError: If the file extension is not supported.
        """
        if not (file or path):
            raise ValueError("No file or path provided.")
        
        extension = splitext(path if path else file.filename)

        if extension not in self.file_types:
            raise ValueError(f"Unsupported file extension: {extension}")
        return self.file_types[extension](file=file, path=path)